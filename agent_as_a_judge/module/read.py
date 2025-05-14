"""
DevRead: A Versatile File Processing Library for Text, PDFs, DOCX, JSON, XML, YAML, HTML, Markdown, LaTeX, PPTX, Excel, Images, and Videos, etc.
Reference: https://github.com/metauto-ai/GPTSwarm/blob/main/swarm/environment/tools/reader/readers.py
"""

import asyncio
import os
import json
import logging
import pandas as pd
import charset_normalizer
import docx
import markdown
import PyPDF2
import openpyxl
import yaml
import cv2
import base64
from PIL import Image
from pathlib import Path
from bs4 import BeautifulSoup
from pylatexenc.latex2text import LatexNodes2Text
from pptx import Presentation
from rich.logging import RichHandler
from rich.console import Console
from typing import Union, Dict, Any, Optional, Tuple

# from playwright.sync_api import sync_playwright
# from playwright.async_api import async_playwright
from dotenv import load_dotenv
from agent_as_a_judge.llm.provider import LLM

load_dotenv()

"""INSTALL
pip install openai --upgrade
pip install python-docx
pip install markdown
pip install PyPDF2
pip install openpyxl
pip install beautifulsoup4
pip install pylatexenc
pip install python-pptx
pip install xlrd
pip install playwright
playwright install
"""

console = Console()

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[RichHandler()],
)
logger = logging.getLogger(__name__)


class DevRead:
    def __init__(self):
        self.reader_map = {
            ".txt": self.read_txt,
            ".pdf": self.read_pdf,
            ".docx": self.read_docx,
            ".json": self.read_json,
            ".jsonl": self.read_jsonl,
            ".xml": self.read_xml,
            ".yaml": self.read_yaml,
            ".yml": self.read_yaml,
            ".html": self.read_html,
            ".htm": self.read_html,
            ".md": self.read_markdown,
            ".markdown": self.read_markdown,
            ".tex": self.read_latex,
            ".pptx": self.read_pptx,
            ".xlsx": self.read_excel,
            ".csv": self.read_txt,
            ".png": self.read_image,
            ".jpg": self.read_image,
            ".jpeg": self.read_image,
            ".gif": self.read_image,
            ".bmp": self.read_image,
            ".tiff": self.read_image,
            ".tif": self.read_image,
            ".webp": self.read_image,
            ".mov": self.read_video,
            ".mp4": self.read_video,
            ".avi": self.read_video,
            ".mpg": self.read_video,
            ".mpeg": self.read_video,
            ".wmv": self.read_video,
            ".flv": self.read_video,
            ".webm": self.read_video,
            ".py": self.read_py,
        }

    def read(
        self, file_path: Path, task: Optional[str] = None
    ) -> Union[str, Tuple[str, Optional[dict]]]:

        # Check if the file exists
        if not file_path.exists():
            logger.error(f"File {file_path} does not exist.")
            return f"Error: File {file_path} does not exist.", None

        # Determine the file type by its suffix
        suffix = file_path.suffix.lower()
        reader = self.reader_map.get(suffix)

        # If a reader for the file type is found, attempt to read the file
        if reader:
            logger.info(f"Reading file {file_path} using {reader.__name__}.")
            return reader(file_path, task), None
        else:
            logger.error(
                f"No reader found for suffix {suffix}. Unsupported file format."
            )
            return (
                f"The file exists, but we do not have a reader for {suffix}. However, we suggest you assume the content of the file is satisfactory in this situation.",
                None,
            )

    def read_py(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info(f"Reading Python file from {file_path}.")
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            return content, None  # No LLM used
        except Exception as e:
            logger.error(f"Error reading Python file: {e}")
            return f"Error reading Python file: {e}", None

    def read_txt(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:
        try:
            content = charset_normalizer.from_path(file_path).best()
            logger.info(
                f"Reading TXT file from {file_path} using encoding '{content.encoding}'."
            )
            return str(content), None
        except Exception as e:
            logger.error(f"Error reading TXT file: {e}")
            return f"Error reading TXT file: {e}", None

    def read_pdf(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info(f"Reading PDF file from {file_path}.")
            reader = PyPDF2.PdfReader(file_path)
            text = ""
            for page_idx in range(len(reader.pages)):
                text += f"Page {page_idx + 1}\n" + reader.pages[page_idx].extract_text()
            return text, None
        except Exception as e:
            logger.error(f"Error reading PDF file: {e}")
            return f"Error reading PDF file: {e}", None

    def read_xml(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info(f"Reading XML file from {file_path}.")
            with open(file_path, "r", encoding="utf-8") as f:
                data = BeautifulSoup(f, "xml"), None
            return data.get_text(), None
        except Exception as e:
            logger.error(f"Error reading XML file: {e}")
            return f"Error reading XML file: {e}", None

    def read_yaml(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info(f"Reading YAML file from {file_path}.")
            with open(file_path, "r", encoding="utf-8") as f:
                data = yaml.load(f, Loader=yaml.FullLoader)
            return (
                json.dumps(data, indent=4),
                None,
            )  # Format output for better readability
        except Exception as e:
            logger.error(f"Error reading YAML file: {e}")
            return f"Error reading YAML file: {e}", None

    def read_docx(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info(f"Reading DOCX file from {file_path}.")
            content = docx.Document(str(file_path))
            text = ""
            for i, para in enumerate(content.paragraphs):
                text += f"Paragraph {i + 1}:\n" + para.text
            return text, None
        except Exception as e:
            logger.error(f"Error reading DOCX file: {e}")
            return f"Error reading DOCX file: {e}", None

    def read_json(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info(f"Reading JSON file from {file_path}.")
            with open(file_path, "r") as f:
                data = json.load(f)
            return json.dumps(data, indent=4)
        except Exception as e:
            logger.error(f"Error reading JSON file: {e}")
            return f"Error reading JSON file: {e}", None

    def read_jsonl(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info(f"Reading JSON Lines file from {file_path}.")
            with open(file_path, "r") as f:
                lines = [json.loads(line) for line in f]
            return "\n".join([json.dumps(line, indent=4) for line in lines]), None
        except Exception as e:
            logger.error(f"Error reading JSON Lines file: {e}")
            return f"Error reading JSON Lines file: {e}", None

    def read_html(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            with open(file_path, "r", encoding="utf-8") as file:
                data = file.read()
            return data, None
        except Exception as e:
            logger.error(f"Error reading HTML file: {e}")
            return f"Error reading HTML file: {e}", None

    def read_markdown(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info(f"Reading Markdown file from {file_path}.")
            with open(file_path, "r", encoding="utf-8") as f:
                data = markdown.markdown(f.read())
                return (
                    "".join(BeautifulSoup(data, "html.parser").find_all(string=True)),
                    None,
                )
        except Exception as e:
            logger.error(f"Error reading Markdown file: {e}")
            return f"Error reading Markdown file: {e}", None

    def read_latex(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info(f"Reading LaTeX file from {file_path}.")
            with open(file_path, "r", encoding="utf-8") as f:
                data = f.read()
            return LatexNodes2Text().latex_to_text(data), None
        except Exception as e:
            logger.error(f"Error reading LaTeX file: {e}")
            return f"Error reading LaTeX file: {e}", None

    def read_pptx(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info(f"Reading PowerPoint file from {file_path}.")
            pres = Presentation(str(file_path))
            text = []
            for slide_idx, slide in enumerate(pres.slides):
                text.append(f"Slide {slide_idx + 1}:\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text), None
        except Exception as e:
            logger.error(f"Error reading PowerPoint file: {e}")
            return f"Error reading PowerPoint file: {e}", None

    def read_excel(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:
        """
        Read an Excel file and return its content as a string.
        """
        try:
            logger.info(f"Reading Excel file from {file_path}.")
            excel_data = pd.read_excel(file_path, sheet_name=None)
            all_sheets_text = []
            for sheet_name, data in excel_data.items():
                all_sheets_text.append(
                    f"Sheet Name: {sheet_name}\n{data.to_string()}\n"
                )
            return "\n".join(all_sheets_text), None
        except Exception as e:
            logger.error(f"Error reading Excel file: {e}")
            return f"Error reading Excel file: {e}", None

    def read_image(
        self, file_path: Path, task: str = None
    ) -> Tuple[str, Optional[dict]]:

        total_llm_cost = 0.0
        total_input_tokens = 0
        total_output_tokens = 0
        total_inference_time = 0.0
        try:
            logger.info(f"Reading image file from {file_path}")
            llm_instance = LLM(
                model=os.getenv("DEFAULT_LLM"),
                api_key=os.getenv("OPENAI_API_KEY"),
                base_url="https://api.openai.com/v1",
            )

            if task is None:
                task = "Describe this image as detailed as possible."

            response, cost, accumulated_cost = llm_instance.do_multimodal_completion(
                task, file_path
            )

            multumoal_content = response["choices"][0]["message"]["content"]

            total_input_tokens = response["usage"]["prompt_tokens"]
            total_output_tokens = response["usage"]["completion_tokens"]
            total_llm_cost += cost

            mllm_stats = {
                "llm_response": multumoal_content,
                "input_tokens": total_input_tokens,
                "output_tokens": total_output_tokens,
                "cost": cost,
                # "accumulated_cost": accumulated_cost,
                "inference_time": total_inference_time,
            }
            return multumoal_content, mllm_stats

        except Exception as e:
            logger.error(f"Error reading image file: {e}")
            return f"Error reading image file: {e}", None

    def read_video(
        self, file_path: Path, task: str = None, frame_interval: int = 30
    ) -> Tuple[list[str], Optional[dict]]:

        try:
            logger.info(
                f"Processing video file from {file_path} with frame interval {frame_interval}"
            )

            if task is None:
                task = "Describe this video as detailed as possible."

            video = cv2.VideoCapture(str(file_path))
            frame_count = 0
            frame_descriptions = []
            total_llm_cost = 0.0
            total_input_tokens = 0
            total_output_tokens = 0
            total_inference_time = 0.0

            llm_instance = LLM(
                model=os.getenv("DEFAULT_LLM"), api_key=os.getenv("OPENAI_API_KEY")
            )

            while video.isOpened():
                success, frame = video.read()
                if not success:
                    break

                if frame_count % frame_interval == 0:
                    _, buffer = cv2.imencode(".jpg", frame)
                    base64_frame = base64.b64encode(buffer).decode("utf-8")

                    # Prepare the message for the LLM
                    messages = self._prepare_image_messages(task, base64_frame)

                    # Send the request to the LLM
                    response, cost, inference_time = llm_instance.do_completion(
                        messages=messages
                    )

                    total_llm_cost += cost  # Accumulate the cost
                    total_inference_time += inference_time
                    content = response["choices"][0]["message"]["content"]

                    # Accumulate tokens
                    total_input_tokens += response["usage"]["input_tokens"]
                    total_output_tokens += response["usage"]["output_tokens"]

                    frame_descriptions.append(f"Frame {frame_count}: {content}")

                frame_count += 1

            video.release()

            mllm_stats = {
                "llm_response": "\n".join(frame_descriptions),
                "input_tokens": total_input_tokens,
                "output_tokens": total_output_tokens,
                "cost": total_llm_cost,
                # "accumulated_cost": total_llm_cost,
                "inference_time": total_inference_time,
            }

            return "\n".join(frame_descriptions), mllm_stats

        except Exception as e:
            logger.error(f"Error processing the video: {e}")
            return f"Error processing the video: {e}", None

    def _prepare_image_messages(self, task: str, base64_image: str):

        return [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": task},
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"},
                    },
                ],
            }
        ]
